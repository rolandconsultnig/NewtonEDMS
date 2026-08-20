"""Microsoft Office Document Properties, Metadata Synchronization, Template Engine & Exporter.

Handles:
1. OpenXML (Word .docx, Excel .xlsx, PowerPoint .pptx) core and custom properties read/write.
2. Word (.docx) & Excel (.xlsx) template placeholder merging (e.g. {{title}}, {{custom_field}}).
3. Professional Excel (.xlsx) dataset exporter with styles, auto-filter, and column widths.
4. Professional Word (.docx) document dossier / report generator.
"""
from __future__ import annotations

import datetime
import io
import logging
import os
import re
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger("newtonedms.office_meta")


def get_office_properties(file_path: Path) -> dict[str, Any]:
    """Read core metadata and custom properties from an Office document (.docx, .xlsx, .pptx)."""
    ext = file_path.suffix.lower()
    props: dict[str, Any] = {
        "title": "",
        "author": "",
        "subject": "",
        "keywords": "",
        "comments": "",
        "category": "",
        "last_modified_by": "",
        "created": None,
        "modified": None,
        "custom": {},
    }
    
    if not file_path.exists():
        return props
    
    # 1. Word Document (.docx)
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(file_path))
            cp = doc.core_properties
            props["title"] = cp.title or ""
            props["author"] = cp.author or ""
            props["subject"] = cp.subject or ""
            props["keywords"] = cp.keywords or ""
            props["comments"] = cp.comments or ""
            props["category"] = cp.category or ""
            props["last_modified_by"] = cp.last_modified_by or ""
            if cp.created:
                props["created"] = cp.created.isoformat()
            if cp.modified:
                props["modified"] = cp.modified.isoformat()
            
            # Read custom properties if present in part
            try:
                part = doc.part
                if hasattr(part, "custom_properties"):
                    for item in part.custom_properties:
                        props["custom"][item.name] = item.value
            except Exception:
                pass
        except Exception as exc:
            logger.warning("Error reading .docx properties: %s", exc)
            
    # 2. Excel Spreadsheet (.xlsx)
    elif ext in (".xlsx", ".xlsm"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=False, data_only=True)
            props["title"] = wb.properties.title or ""
            props["author"] = wb.properties.creator or ""
            props["subject"] = wb.properties.subject or ""
            props["keywords"] = wb.properties.keywords or ""
            props["comments"] = wb.properties.description or ""
            props["category"] = wb.properties.category or ""
            props["last_modified_by"] = wb.properties.lastModifiedBy or ""
            if wb.properties.created:
                props["created"] = wb.properties.created.isoformat()
            if wb.properties.modified:
                props["modified"] = wb.properties.modified.isoformat()
            
            # Custom properties in openpyxl
            if hasattr(wb, "custom_doc_props") and wb.custom_doc_props:
                for p in wb.custom_doc_props:
                    props["custom"][p.name] = p.value
            wb.close()
        except Exception as exc:
            logger.warning("Error reading .xlsx properties: %s", exc)
            
    # 3. PowerPoint Presentation (.pptx)
    elif ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(str(file_path))
            cp = prs.core_properties
            props["title"] = cp.title or ""
            props["author"] = cp.author or ""
            props["subject"] = cp.subject or ""
            props["keywords"] = cp.keywords or ""
            props["comments"] = cp.comments or ""
            props["category"] = cp.category or ""
            props["last_modified_by"] = cp.last_modified_by or ""
            if cp.created:
                props["created"] = cp.created.isoformat()
            if cp.modified:
                props["modified"] = cp.modified.isoformat()
        except Exception as exc:
            logger.warning("Error reading .pptx properties: %s", exc)
            
    return props


def update_office_properties(file_path: Path, new_props: dict[str, Any]) -> bool:
    """Write/synchronize core metadata and custom properties into an Office document."""
    ext = file_path.suffix.lower()
    if not file_path.exists():
        return False
    
    # 1. Word Document (.docx)
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(file_path))
            cp = doc.core_properties
            if "title" in new_props:
                cp.title = str(new_props["title"])
            if "author" in new_props:
                cp.author = str(new_props["author"])
            if "subject" in new_props:
                cp.subject = str(new_props["subject"])
            if "keywords" in new_props:
                cp.keywords = str(new_props["keywords"])
            if "comments" in new_props:
                cp.comments = str(new_props["comments"])
            if "category" in new_props:
                cp.category = str(new_props["category"])
            cp.last_modified_by = "NewtonEDMS"
            cp.modified = datetime.datetime.now(datetime.timezone.utc)
            doc.save(str(file_path))
            return True
        except Exception as exc:
            logger.warning("Error writing .docx properties: %s", exc)
            return False
            
    # 2. Excel Spreadsheet (.xlsx)
    elif ext in (".xlsx", ".xlsm"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path))
            if "title" in new_props:
                wb.properties.title = str(new_props["title"])
            if "author" in new_props:
                wb.properties.creator = str(new_props["author"])
            if "subject" in new_props:
                wb.properties.subject = str(new_props["subject"])
            if "keywords" in new_props:
                wb.properties.keywords = str(new_props["keywords"])
            if "comments" in new_props:
                wb.properties.description = str(new_props["comments"])
            if "category" in new_props:
                wb.properties.category = str(new_props["category"])
            wb.properties.lastModifiedBy = "NewtonEDMS"
            wb.properties.modified = datetime.datetime.now(datetime.timezone.utc)
            wb.save(str(file_path))
            wb.close()
            return True
        except Exception as exc:
            logger.warning("Error writing .xlsx properties: %s", exc)
            return False
            
    # 3. PowerPoint Presentation (.pptx)
    elif ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(str(file_path))
            cp = prs.core_properties
            if "title" in new_props:
                cp.title = str(new_props["title"])
            if "author" in new_props:
                cp.author = str(new_props["author"])
            if "subject" in new_props:
                cp.subject = str(new_props["subject"])
            if "keywords" in new_props:
                cp.keywords = str(new_props["keywords"])
            if "comments" in new_props:
                cp.comments = str(new_props["comments"])
            if "category" in new_props:
                cp.category = str(new_props["category"])
            cp.last_modified_by = "NewtonEDMS"
            cp.modified = datetime.datetime.now(datetime.timezone.utc)
            prs.save(str(file_path))
            return True
        except Exception as exc:
            logger.warning("Error writing .pptx properties: %s", exc)
            return False
            
    return False


def merge_word_template(template_path: Path, output_path: Path, context: dict[str, Any]) -> Path:
    """Merge placeholder values (e.g. {{name}}, {{title}}, {{date}}) into a Word .docx template."""
    import docx
    
    doc = docx.Document(str(template_path))
    
    # Helper to replace placeholders in a paragraph
    def _replace_in_paragraph(p):
        full_text = p.text
        for key, val in context.items():
            pattern = r"\{\{\s*" + re.escape(str(key)) + r"\s*\}\}"
            if re.search(pattern, full_text):
                full_text = re.sub(pattern, str(val if val is not None else ""), full_text)
        if full_text != p.text:
            p.text = full_text

    # 1. Paragraphs
    for p in doc.paragraphs:
        _replace_in_paragraph(p)
        
    # 2. Tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _replace_in_paragraph(p)
                    
    # 3. Headers and footers
    for section in doc.sections:
        for hp in section.header.paragraphs:
            _replace_in_paragraph(hp)
        for fp in section.footer.paragraphs:
            _replace_in_paragraph(fp)
            
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return output_path


def merge_excel_template(template_path: Path, output_path: Path, context: dict[str, Any]) -> Path:
    """Merge placeholder values into an Excel .xlsx template across all sheets and cells."""
    import openpyxl
    
    wb = openpyxl.load_workbook(str(template_path))
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    val = cell.value
                    for k, v in context.items():
                        pattern = r"\{\{\s*" + re.escape(str(k)) + r"\s*\}\}"
                        if re.search(pattern, val):
                            val = re.sub(pattern, str(v if v is not None else ""), val)
                    cell.value = val
                    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))
    wb.close()
    return output_path


def export_documents_to_excel(documents: list[dict], title: str = "NewtonEDMS Export") -> bytes:
    """Generate a professionally styled Excel workbook (.xlsx) containing document records."""
    import openpyxl
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Documents"
    ws.views.sheetView[0].showGridLines = True
    
    # Brand styles
    title_font = Font(name="Segoe UI", size=16, bold=True, color="1E293B")
    sub_font = Font(name="Segoe UI", size=10, italic=True, color="64748B")
    header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    
    zebra_fill = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")
    data_font = Font(name="Segoe UI", size=10, color="0F172A")
    
    thin_border = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )
    
    # Title block
    ws.merge_cells("A1:H1")
    ws["A1"] = f"NewtonEDMS — {title}"
    ws["A1"].font = title_font
    ws["A1"].alignment = Alignment(vertical="center")
    
    ws.merge_cells("A2:H2")
    ws["A2"] = f"Exported on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')} UTC | Total records: {len(documents)}"
    ws["A2"].font = sub_font
    ws["A2"].alignment = Alignment(vertical="center")
    
    ws.row_dimensions[1].height = 28
    ws.row_dimensions[2].height = 18
    ws.row_dimensions[4].height = 24
    
    # Headers
    headers = [
        "Doc ID",
        "Document Name",
        "Folder",
        "Version",
        "Status",
        "Size (Bytes)",
        "Tags / Categories",
        "Last Modified",
    ]
    
    start_row = 4
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=start_row, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center" if col_idx in (1, 4, 5) else "left", vertical="center")
        cell.border = thin_border
        
    # Populate rows
    for row_idx, doc in enumerate(documents, start_row + 1):
        ws.row_dimensions[row_idx].height = 20
        tags_str = ", ".join(doc.get("tags") or []) if isinstance(doc.get("tags"), list) else str(doc.get("tags") or "")
        
        row_values = [
            doc.get("id"),
            doc.get("name") or "Untitled",
            doc.get("folder_name") or f"Folder #{doc.get('folder_id', 0)}",
            doc.get("version") or "1.0",
            doc.get("status") or "draft",
            doc.get("size") or 0,
            tags_str,
            str(doc.get("updated_at") or doc.get("created_at") or "")[:19],
        ]
        
        is_even = (row_idx % 2 == 0)
        for col_idx, val in enumerate(row_values, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.font = data_font
            cell.border = thin_border
            if is_even:
                cell.fill = zebra_fill
            if col_idx in (1, 4, 5):
                cell.alignment = Alignment(horizontal="center", vertical="center")
            elif col_idx == 6:
                cell.alignment = Alignment(horizontal="right", vertical="center")
                cell.number_format = "#,##0"
            else:
                cell.alignment = Alignment(horizontal="left", vertical="center")
                
    # Auto-filter and Column Widths
    end_row = start_row + max(len(documents), 1)
    ws.auto_filter.ref = f"A{start_row}:H{end_row}"
    
    for col in ws.columns:
        max_len = 0
        col_letter = get_column_letter(col[0].column)
        for cell in col:
            if cell.row >= start_row and cell.value:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)
        
    buf = io.BytesIO()
    wb.save(buf)
    wb.close()
    return buf.getvalue()


def export_documents_to_word(documents: list[dict], title: str = "NewtonEDMS Document Dossier") -> bytes:
    """Generate a formatted Word document (.docx) summary report."""
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    
    doc = docx.Document()
    
    # Title
    h1 = doc.add_heading(title, level=0)
    h1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub.add_run(f"Generated by NewtonEDMS on {datetime.datetime.now().strftime('%B %d, %Y at %H:%M:%S UTC')}")
    run.font.italic = True
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(100, 116, 139)
    
    doc.add_paragraph()
    
    # Overview table
    doc.add_heading("Repository Document Summary", level=2)
    t = doc.add_table(rows=1, cols=5)
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.autofit = False
    
    headers = ["ID", "Document Name", "Folder", "Status", "Version"]
    hdr_cells = t.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        hdr_cells[i].paragraphs[0].runs[0].font.bold = True
        hdr_cells[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        # Background color
        tcPr = hdr_cells[i]._element.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), "1E3A8A")
        tcPr.append(shd)
        
    for d in documents:
        row_cells = t.add_row().cells
        row_cells[0].text = str(d.get("id", ""))
        row_cells[1].text = str(d.get("name", ""))
        row_cells[2].text = str(d.get("folder_name", "") or f"#{d.get('folder_id', '')}")
        row_cells[3].text = str(d.get("status", "draft"))
        row_cells[4].text = str(d.get("version", "1.0"))
        
    doc.add_paragraph()
    
    # Detailed entries
    if documents:
        doc.add_heading("Detailed Item Records", level=2)
        for d in documents[:50]: # limit to 50 detailed entries per dossier
            doc.add_heading(f"#{d.get('id')} — {d.get('name')}", level=3)
            p = doc.add_paragraph()
            p.add_run("Status: ").bold = True
            p.add_run(f"{d.get('status', 'draft')}    |    ")
            p.add_run("Version: ").bold = True
            p.add_run(f"{d.get('version', '1.0')}    |    ")
            p.add_run("Size: ").bold = True
            p.add_run(f"{d.get('size', 0):,} bytes\n")
            
            tags = d.get("tags") or []
            if tags:
                p.add_run("Tags: ").bold = True
                p.add_run(f"{', '.join(tags) if isinstance(tags, list) else tags}\n")
            
            notes = d.get("notes") or d.get("description")
            if notes:
                p.add_run("Notes: ").bold = True
                p.add_run(f"{notes}\n")
                
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
