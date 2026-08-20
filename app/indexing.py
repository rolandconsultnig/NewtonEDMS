"""Full-text extraction and Whoosh search indexing.

Text is extracted from common office formats, PDFs and images (OCR).
The Whoosh index lives under ``storage/whoosh_index``.
"""
from __future__ import annotations

import logging
import zipfile
from pathlib import Path

from app import database
from app.config import settings

logger = logging.getLogger("newtonedms.indexing")

def _index_dir() -> Path:
    # Resolved per call so runtime storage overrides (and tests) are honoured.
    return database.STORAGE_DIR / "whoosh_index"


def _imports():
    """Lazy imports so the app starts even when optional libraries are absent."""
    mods = {}
    for name in ["whoosh", "whoosh.index", "whoosh.fields", "whoosh.qparser", "whoosh.writing"]:
        try:
            __import__(name)
            mods[name] = __import__(name)
        except ImportError:
            pass
    return mods


def _ensure_index():
    from whoosh import index
    from whoosh.fields import ID, NUMERIC, TEXT, Schema

    _index_dir().mkdir(parents=True, exist_ok=True)
    if not index.exists_in(_index_dir()):
        schema = Schema(
            doc_id=ID(stored=True, unique=True),
            title=TEXT(stored=True),
            tags=TEXT(stored=True),
            content=TEXT(stored=False),
            barcodes=TEXT(stored=True),
            size=NUMERIC(stored=True, sortable=True),
        )
        return index.create_in(_index_dir(), schema)
    return index.open_dir(_index_dir())


def _extract_text(path: Path, mime: str) -> tuple[str, str]:
    """Return (text, barcodes). Best-effort extraction with a size guard."""
    text = ""
    barcodes = ""
    try:
        if path.stat().st_size > settings.max_extract_bytes:
            # Multi-GB files must not be read into memory in the request path.
            return text, barcodes
    except OSError:
        return text, barcodes
    ext = path.suffix.lower()

    # Plain text / code
    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log", ".py", ".js", ".css", ".sql"):
        try:
            # Stream with a byte budget rather than reading the whole file.
            budget = settings.max_extract_bytes
            chunks: list[str] = []
            with path.open("r", encoding="utf-8", errors="ignore") as fh:
                while budget > 0:
                    chunk = fh.read(min(1024 * 1024, budget))
                    if not chunk:
                        break
                    budget -= len(chunk.encode("utf-8", errors="ignore"))
                    chunks.append(chunk)
                text = "".join(chunks)
        except Exception:
            pass
        return text, barcodes

    # PDF
    if ext == ".pdf":
        try:
            import pdfplumber

            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception:
            pass

    # Word
    if ext in (".docx", ".doc"):
        try:
            import docx

            doc = docx.Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            pass

    # Excel
    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    text += " ".join(str(c) for c in row if c is not None) + "\n"
        except Exception:
            pass

    # PowerPoint
    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation

            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception:
            pass

    # Images: OCR + barcode
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"):
        try:
            from PIL import Image

            img = Image.open(path)
            # OCR
            try:
                import pytesseract

                text = pytesseract.image_to_string(img)
            except Exception:
                pass
            # Barcodes
            try:
                from pyzbar.pyzbar import decode

                codes = decode(img)
                barcodes = " ".join(c.data.decode("utf-8", errors="ignore") for c in codes if c.data)
            except Exception:
                pass
        except Exception:
            pass

    # Zip archives: index filenames
    if ext == ".zip":
        try:
            with zipfile.ZipFile(path) as zf:
                text = " ".join(zf.namelist())
        except Exception:
            pass

    return text, barcodes


def index_document(
    doc_id: int,
    title: str,
    tags: str,
    file_path: str,
    size: int,
    content_override: str | None = None,
) -> None:
    """Add or update a document in the full-text index."""
    try:
        from whoosh.writing import AsyncWriter

        ix = _ensure_index()
        barcodes = ""
        if content_override is not None:
            text = content_override
        elif file_path:
            text, barcodes = _extract_text(Path(file_path), "")
        else:
            text = ""
        writer = AsyncWriter(ix)
        writer.update_document(
            doc_id=str(doc_id),
            title=title or "",
            tags=tags or "",
            content=text,
            barcodes=barcodes,
            size=size,
        )
        writer.commit()
    except Exception:
        # Indexing is best-effort; the app continues to work if libraries are missing.
        logger.warning("index_document failed for doc %s", doc_id, exc_info=True)


def remove_document(doc_id: int) -> None:
    """Remove a document from the full-text index."""
    try:
        from whoosh import index

        if not _index_dir().exists():
            return
        ix = index.open_dir(_index_dir())
        writer = ix.writer()
        writer.delete_by_term("doc_id", str(doc_id))
        writer.commit()
    except Exception:
        logger.warning("remove_document failed for doc %s", doc_id, exc_info=True)


def search_documents(query: str, limit: int = 100) -> list[int]:
    """Return the ids of documents matching the full-text query."""
    try:
        from whoosh import index
        from whoosh.qparser import MultifieldParser

        if not query or not _index_dir().exists():
            return []
        if any(ch in query for ch in ("*", "?", "~", "[", "]", ":")):
            # Reject wildcard/fuzzy/range/field syntax: expensive full-index scans.
            return []
        ix = index.open_dir(_index_dir())
        with ix.searcher() as searcher:
            parser = MultifieldParser(["title", "tags", "content", "barcodes"], ix.schema)
            q = parser.parse(query)
            results = searcher.search(q, limit=limit)
            return [int(r["doc_id"]) for r in results]
    except Exception:
        logger.warning("search_documents failed", exc_info=True)
        return []
